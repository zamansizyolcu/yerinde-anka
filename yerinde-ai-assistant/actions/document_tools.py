"""
Belge analizi: PDF, Word (docx), PowerPoint (pptx), Excel (xlsx) metni
çıkarır; resim dosyalarını Gemini Vision ile analiz eder. Ayrıca PDF'leri
sesli kitap gibi okuyabilir (mevcut TTS motorlarından biriyle, tamamen
yerel çalışır).
"""

from __future__ import annotations

import mimetypes
from pathlib import Path

from app_config import get_app_config_value

MAX_CHARS = 6000
TTS_READ_LIMIT = 4000  # tek seferde sesli okunacak azami karakter

DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".xlsm", ".doc", ".ppt", ".xls"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp"}
# Doğrudan metin olarak okunabilen kod/metin dosyaları (ör. yüklenen .py
# dosyalarını analiz ettirmek / hatalarını buldurmak için)
TEXT_EXTENSIONS = {".py", ".txt", ".md", ".json", ".js", ".ts", ".html", ".css",
                    ".c", ".cpp", ".h", ".java", ".cs", ".sh", ".yaml", ".yml"}


def _extract_plain_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        return f"Dosya okunamadı: {e}"


def _resolve_path(file_path: str) -> Path | None:
    candidate = (file_path or "").strip()
    if not candidate:
        candidate = str(get_app_config_value("last_uploaded_file", "") or "").strip()
    if not candidate:
        return None
    return Path(candidate)


def _extract_pdf_text(path: Path) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(str(path)) as pdf:
            parts = [page.extract_text() or "" for page in pdf.pages]
        return "\n".join(parts).strip()
    except ImportError:
        pass
    except Exception as e:
        return f"PDF okunamadı: {e}"

    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(path))
        parts = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(parts).strip()
    except ImportError:
        return "PDF okumak için 'pip install pdfplumber' (veya 'pip install PyPDF2') gerekli."
    except Exception as e:
        return f"PDF okunamadı: {e}"


def _extract_docx_text(path: Path) -> str:
    try:
        import docx
    except ImportError:
        return "Word (.docx) dosyalarını okumak için 'pip install python-docx' gerekli."
    try:
        d = docx.Document(str(path))
        return "\n".join(p.text for p in d.paragraphs if p.text.strip())
    except Exception as e:
        return f"Word dosyası okunamadı: {e}"


def _extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "PowerPoint (.pptx) dosyalarını okumak için 'pip install python-pptx' gerekli."
    try:
        prs = Presentation(str(path))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"--- Slayt {i} ---")
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                    lines.append(shape.text_frame.text)
        return "\n".join(lines)
    except Exception as e:
        return f"PowerPoint dosyası okunamadı: {e}"


def _extract_xlsx_text(path: Path) -> str:
    try:
        import openpyxl
    except ImportError:
        return "Excel (.xlsx) dosyalarını okumak için 'pip install openpyxl' gerekli."
    try:
        wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
        lines = []
        for ws in wb.worksheets:
            lines.append(f"--- Sayfa: {ws.title} ---")
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    lines.append(" | ".join(str(c) if c is not None else "" for c in row))
                    row_count += 1
                if row_count >= 200:  # çok büyük sayfalarda kesme
                    lines.append("... (devamı kesildi)")
                    break
        return "\n".join(lines)
    except Exception as e:
        return f"Excel dosyası okunamadı: {e}"


def extract_document_text(file_path: str) -> str:
    """Bir belgenin (pdf/docx/pptx/xlsx) metnini çıkarır (uzunsa kısaltır)."""
    path = _resolve_path(file_path)
    if not path:
        return "Analiz edilecek dosya belirtilmedi. Önce 'DOSYA YÜKLE' düğmesiyle bir dosya seç ya da tam yolunu söyle."
    if not path.exists():
        return f"Dosya bulunamadı: {path}"

    ext = path.suffix.lower()
    if ext == ".pdf":
        text = _extract_pdf_text(path)
    elif ext == ".docx":
        text = _extract_docx_text(path)
    elif ext == ".pptx":
        text = _extract_pptx_text(path)
    elif ext in (".xlsx", ".xlsm"):
        text = _extract_xlsx_text(path)
    elif ext in (".doc", ".ppt", ".xls"):
        return (f"'{ext}' eski Office formatı doğrudan okunamıyor. Dosyayı "
                f".docx/.pptx/.xlsx olarak farklı kaydedip tekrar dener misin?")
    elif ext in TEXT_EXTENSIONS:
        text = _extract_plain_text(path)
    else:
        return f"Desteklenmeyen dosya türü: {ext}"

    if not text:
        return f"'{path.name}' içinden metin çıkarılamadı (taranmış görsel bir PDF olabilir)."
    if len(text) > MAX_CHARS:
        text = text[:MAX_CHARS] + f"\n... ('{path.name}' kısaltıldı, toplam {len(text)} karakter)"
    return f"[{path.name} içeriği]\n{text}"


def analyze_image_file(file_path: str, query: str = "") -> str:
    """Bir resim dosyasını Gemini Vision ile analiz eder (internet gerektirir)."""
    path = _resolve_path(file_path)
    if not path:
        return "Analiz edilecek resim belirtilmedi. Önce 'DOSYA YÜKLE' düğmesiyle bir resim seç."
    if not path.exists():
        return f"Dosya bulunamadı: {path}"

    api_key = str(get_app_config_value("gemini_api_key", "") or "").strip()
    if not api_key:
        return ("Resim analizi Gemini Vision kullanır ve bir API anahtarı gerektirir "
                "(çevrimdışı Ollama modunda resim analizi desteklenmiyor).")

    try:
        from google import genai
        from google.genai import types
    except ImportError:
        return "google-genai paketi yüklü değil."

    mime = mimetypes.guess_type(str(path))[0] or "image/jpeg"
    prompt = query.strip() if query and query.strip() else "Bu resimde ne var? Türkçe, kısa ve net anlat."

    try:
        client = genai.Client(api_key=api_key)
        image_bytes = path.read_bytes()
        for model in ("models/gemini-2.0-flash", "models/gemini-2.5-flash-lite", "models/gemini-2.5-flash"):
            try:
                resp = client.models.generate_content(
                    model=model,
                    contents=[types.Part.from_bytes(data=image_bytes, mime_type=mime), prompt],
                )
                if resp and resp.text:
                    return resp.text.strip()
            except Exception:
                continue
        return "Resim analiz edilemedi (hiçbir model yanıt vermedi)."
    except Exception as e:
        return f"Resim analiz hatası: {e}"


def analyze_document(file_path: str = "", query: str = "") -> str:
    """
    Bir belgeyi ('yükle' düğmesiyle seçilmiş ya da yolu verilmiş) türüne göre
    ele alır: pdf/docx/pptx/xlsx için metin çıkarır, resimler için Gemini
    Vision ile görsel analiz yapar.
    """
    path = _resolve_path(file_path)
    if not path:
        return "Analiz edilecek dosya belirtilmedi. Önce 'DOSYA YÜKLE' düğmesiyle bir dosya seç."
    if not path.exists():
        return f"Dosya bulunamadı: {path}"

    ext = path.suffix.lower()
    if ext in IMAGE_EXTENSIONS:
        return analyze_image_file(str(path), query)
    if ext in DOCUMENT_EXTENSIONS or ext in TEXT_EXTENSIONS:
        return extract_document_text(str(path))
    return f"Desteklenmeyen dosya türü: {ext}"


def read_document_aloud(file_path: str = "") -> str:
    """
    Bir PDF/Word/PowerPoint dosyasının metnini çıkarıp sesli kitap gibi okur.
    Mevcut TTS tercihini kullanır (Piper / SAPI / espeak-ng — tamamen yerel).
    Çok uzun belgelerde ilk ~4000 karakteri okur ve kalanını bildirir.
    """
    from actions.offline_tts import speak_text_offline

    path = _resolve_path(file_path)
    if not path:
        return "Sesli okunacak dosya belirtilmedi. Önce 'DOSYA YÜKLE' düğmesiyle bir dosya seç."
    if not path.exists():
        return f"Dosya bulunamadı: {path}"

    ext = path.suffix.lower()
    if ext == ".pdf":
        text = _extract_pdf_text(path)
    elif ext == ".docx":
        text = _extract_docx_text(path)
    elif ext == ".pptx":
        text = _extract_pptx_text(path)
    else:
        return f"Sesli okuma yalnızca PDF/Word/PowerPoint için destekleniyor ('{ext}' değil)."

    if not text or not text.strip():
        return f"'{path.name}' içinden okunacak metin bulunamadı."

    remaining = len(text) - TTS_READ_LIMIT
    chunk = text[:TTS_READ_LIMIT]

    speak_text_offline(chunk, blocking=True)

    if remaining > 0:
        return f"'{path.name}' dosyasının ilk bölümünü okudum (~{remaining} karakter kaldı — 'devamını oku' diyebilirsin)."
    return f"'{path.name}' dosyasını tamamen okudum."
